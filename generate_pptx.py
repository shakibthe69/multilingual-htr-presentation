import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set slide dimensions to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    c_navy = RGBColor(10, 15, 29)
    c_card = RGBColor(19, 26, 48)
    c_white = RGBColor(248, 250, 252)
    c_secondary = RGBColor(148, 163, 184)
    c_cyan = RGBColor(0, 240, 255)
    c_blue = RGBColor(56, 189, 248)
    c_purple = RGBColor(168, 85, 247)
    c_red = RGBColor(248, 113, 113)
    c_green = RGBColor(16, 185, 129)

    # Helper function to set slide background
    def set_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = c_navy
        
        # Add "Built by Shakib" author tag to bottom right corner
        author_box = slide.shapes.add_textbox(Inches(10.5), Inches(6.9), Inches(2.25), Inches(0.4))
        tf = author_box.text_frame
        p = tf.paragraphs[0]
        p.text = "BUILT BY SHAKIB"
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(8.5)
        p.font.color.rgb = c_secondary
        p.font.bold = True
        p.font.name = 'Trebuchet MS'

    # Helper to add standard slide header
    def add_header(slide, title_text, category_text=None):
        if category_text:
            cat_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.4))
            tf_cat = cat_box.text_frame
            tf_cat.word_wrap = True
            p_cat = tf_cat.paragraphs[0]
            p_cat.text = category_text.upper()
            p_cat.font.size = Pt(10)
            p_cat.font.bold = True
            p_cat.font.color.rgb = c_blue
            p_cat.font.name = 'Trebuchet MS'
            
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.7), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = c_white
        p.font.name = 'Trebuchet MS'
        
        # Add a subtle line under the header
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.5), Inches(12.133), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(30, 41, 59)
        line.line.fill.background()

    blank_layout = prs.slide_layouts[6]

    # =========================================================================
    # SLIDE 1: TITLE SLIDE
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_bg(slide1)
    
    # Left Column (Text Info)
    left_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(6.5), Inches(5.0))
    tf1 = left_box.text_frame
    tf1.word_wrap = True
    
    # Badge
    p_badge = tf1.paragraphs[0]
    p_badge.text = "RESEARCH PROJECT • CSE"
    p_badge.font.size = Pt(11)
    p_badge.font.bold = True
    p_badge.font.color.rgb = c_purple
    p_badge.font.name = 'Trebuchet MS'
    
    # Main Title
    p_title = tf1.add_paragraph()
    p_title.text = "\nMultilingual Handwritten\nText Recognition"
    p_title.font.size = Pt(40)
    p_title.font.bold = True
    p_title.font.color.rgb = c_white
    p_title.font.name = 'Trebuchet MS'
    
    # Subtitle
    p_sub = tf1.add_paragraph()
    p_sub.text = "Bangla • English • Digits • Symbols\n"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = c_cyan
    p_sub.font.name = 'Trebuchet MS'
    
    # Tagline
    p_tag = tf1.add_paragraph()
    p_tag.text = "Converting handwritten images into editable digital text."
    p_tag.font.size = Pt(13)
    p_tag.font.italic = True
    p_tag.font.color.rgb = c_secondary
    p_tag.font.name = 'Arial'
    
    # Meta Info (Bottom Left)
    meta_box = slide1.shapes.add_textbox(Inches(0.8), Inches(5.1), Inches(6.5), Inches(2.0))
    tf_meta = meta_box.text_frame
    tf_meta.word_wrap = True
    p_meta = tf_meta.paragraphs[0]
    p_meta.text = "PREPARED BY:\nMD SHAKIB AHMED (112230862)  •  Md Fahad Molla (112230890)\nMuhammad Yeakub (112231003)  •  Md. Imjam Hosen (112230335)\n\nDepartment of Computer Science & Engineering\nUnited International University"
    p_meta.font.size = Pt(10)
    p_meta.font.color.rgb = c_secondary
    p_meta.font.name = 'Arial'

    # Right Column (Visual Pipeline Card)
    card_bg = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1.8), Inches(4.7), Inches(4.5))
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = c_card
    card_bg.line.color.rgb = RGBColor(30, 41, 59)
    
    # Pipeline components inside Right Column
    pipe_box = slide1.shapes.add_textbox(Inches(8.0), Inches(2.0), Inches(4.3), Inches(4.1))
    tf_pipe = pipe_box.text_frame
    tf_pipe.word_wrap = True
    
    p_pipe_title = tf_pipe.paragraphs[0]
    p_pipe_title.text = "TRANSFORMATION PIPELINE"
    p_pipe_title.alignment = PP_ALIGN.CENTER
    p_pipe_title.font.size = Pt(12)
    p_pipe_title.font.bold = True
    p_pipe_title.font.color.rgb = c_blue
    
    p_pipe_step1 = tf_pipe.add_paragraph()
    p_pipe_step1.text = "\n[ Handwritten Input ]\nবাংলা  •  Hello  •  2026  •  @ + ="
    p_pipe_step1.alignment = PP_ALIGN.CENTER
    p_pipe_step1.font.size = Pt(12)
    p_pipe_step1.font.color.rgb = c_white
    
    p_arrow1 = tf_pipe.add_paragraph()
    p_arrow1.text = "↓"
    p_arrow1.alignment = PP_ALIGN.CENTER
    p_arrow1.font.size = Pt(16)
    p_arrow1.font.color.rgb = c_purple
    
    p_pipe_step2 = tf_pipe.add_paragraph()
    p_pipe_step2.text = "[ AI Vision Model ]"
    p_pipe_step2.alignment = PP_ALIGN.CENTER
    p_pipe_step2.font.size = Pt(13)
    p_pipe_step2.font.bold = True
    p_pipe_step2.font.color.rgb = c_cyan
    
    p_arrow2 = tf_pipe.add_paragraph()
    p_arrow2.text = "↓"
    p_arrow2.alignment = PP_ALIGN.CENTER
    p_arrow2.font.size = Pt(16)
    p_arrow2.font.color.rgb = c_purple
    
    p_pipe_step3 = tf_pipe.add_paragraph()
    p_pipe_step3.text = "[ Clean Digital Text ]\nবাংলা  •  Hello  •  2026  •  @ + ="
    p_pipe_step3.alignment = PP_ALIGN.CENTER
    p_pipe_step3.font.size = Pt(12)
    p_pipe_step3.font.color.rgb = c_green

    # =========================================================================
    # SLIDE 2: PROBLEM & MOTIVATION
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_bg(slide2)
    add_header(slide2, "Why Handwriting Recognition?", "Motivation")
    
    # Left Card (Problem)
    left_card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9), Inches(5.6), Inches(4.2))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = c_card
    left_card.line.color.rgb = c_red
    
    lc_box = slide2.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(3.8))
    tf_lc = lc_box.text_frame
    tf_lc.word_wrap = True
    
    p_lc_title = tf_lc.paragraphs[0]
    p_lc_title.text = "⚠️ Complex Handwriting Challenge"
    p_lc_title.font.size = Pt(18)
    p_lc_title.font.bold = True
    p_lc_title.font.color.rgb = c_red
    
    bullet1 = tf_lc.add_paragraph()
    bullet1.text = "\n• Complex Writing Styles\n  Different writing habits and unique curves per individual."
    bullet1.font.size = Pt(12)
    bullet1.font.color.rgb = c_white
    
    bullet2 = tf_lc.add_paragraph()
    bullet2.text = "\n• Noise & Distortion\n  Smudges, low-contrast scans, paper folds, and background noise."
    bullet2.font.size = Pt(12)
    bullet2.font.color.rgb = c_white
    
    bullet3 = tf_lc.add_paragraph()
    bullet3.text = "\n• Mixed Scripts\n  Multiple languages/scripts present within the same document."
    bullet3.font.size = Pt(12)
    bullet3.font.color.rgb = c_white

    # Right Card (Example)
    right_card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.9), Inches(5.6), Inches(4.2))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = c_card
    right_card.line.color.rgb = c_cyan
    
    rc_box = slide2.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.2), Inches(3.8))
    tf_rc = rc_box.text_frame
    tf_rc.word_wrap = True
    
    p_rc_title = tf_rc.paragraphs[0]
    p_rc_title.text = "📑 Real-World Multilingual Sample"
    p_rc_title.font.size = Pt(18)
    p_rc_title.font.bold = True
    p_rc_title.font.color.rgb = c_cyan
    
    p_sample_lbl = tf_rc.add_paragraph()
    p_sample_lbl.text = "\nDocument Sample Content:"
    p_sample_lbl.font.size = Pt(12)
    p_sample_lbl.font.color.rgb = c_secondary
    
    p_sample_text = tf_rc.add_paragraph()
    p_sample_text.text = ' আমার নাম Shakib\n ID: 202012345\n Total = 500 + 200\n Email: shakib@gmail.com'
    p_sample_text.font.size = Pt(15)
    p_sample_text.font.bold = True
    p_sample_text.font.color.rgb = c_white
    p_sample_text.font.name = 'Consolas'
    
    p_sample_res = tf_rc.add_paragraph()
    p_sample_res.text = "\n→ One unified document containing multiple character types."
    p_sample_res.font.size = Pt(13)
    p_sample_res.font.color.rgb = c_cyan
    p_sample_res.font.bold = True

    # Bottom Alert
    bot_alert = slide2.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.6))
    tf_ba = bot_alert.text_frame
    tf_ba.word_wrap = True
    p_ba = tf_ba.paragraphs[0]
    p_ba.text = "Limitation: Existing optical/handwritten systems are often language- or task-specific."
    p_ba.font.size = Pt(12)
    p_ba.font.bold = True
    p_ba.font.color.rgb = c_red

    # =========================================================================
    # SLIDE 3: OBJECTIVES & SCOPE
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_bg(slide3)
    add_header(slide3, "What Are We Building?", "Objectives & Scope")
    
    # 4 Horizontal Cards
    card_width = Inches(2.7)
    card_height = Inches(1.8)
    card_y = Inches(1.8)
    
    cards_data = [
        ("🇧🇩 Bangla", "Handwritten Bangla script and graphemes", 0.8),
        ("🔤 English", "Uppercase & lowercase alphabet recognition", 3.8),
        ("🔢 Digits", "0-9 English & Bengali digits recognition", 6.8),
        ("🔣 Symbols", "@ # + - = % and other special characters", 9.8)
    ]
    
    for title, desc, x in cards_data:
        c_shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), card_y, card_width, card_height)
        c_shape.fill.solid()
        c_shape.fill.fore_color.rgb = c_card
        c_shape.line.color.rgb = c_blue
        
        c_box = slide3.shapes.add_textbox(Inches(x + 0.1), card_y + Inches(0.1), card_width - Inches(0.2), card_height - Inches(0.2))
        tf_c = c_box.text_frame
        tf_c.word_wrap = True
        
        p_t = tf_c.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = c_cyan
        
        p_d = tf_c.add_paragraph()
        p_d.text = "\n" + desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = c_white

    # Core Goal Text block
    goal_box = slide3.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.5))
    tf_g = goal_box.text_frame
    tf_g.word_wrap = True
    p_g = tf_g.paragraphs[0]
    p_g.text = "Core Goal: One system → Multiple handwritten content types"
    p_g.alignment = PP_ALIGN.CENTER
    p_g.font.size = Pt(14)
    p_g.font.bold = True
    p_g.font.color.rgb = c_cyan

    # Scope Split (In Scope vs Out of Scope)
    # Left Scope Box
    in_shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.5), Inches(5.6), Inches(2.2))
    in_shape.fill.solid()
    in_shape.fill.fore_color.rgb = c_card
    in_shape.line.color.rgb = c_green
    
    in_box = slide3.shapes.add_textbox(Inches(0.9), Inches(4.6), Inches(5.4), Inches(2.0))
    tf_in = in_box.text_frame
    tf_in.word_wrap = True
    p_in_title = tf_in.paragraphs[0]
    p_in_title.text = "IN SCOPE"
    p_in_title.font.size = Pt(13)
    p_in_title.font.bold = True
    p_in_title.font.color.rgb = c_green
    
    p_in_bullets = tf_in.add_paragraph()
    p_in_bullets.text = "• Image-based HTR\n• Mixed script recognition\n• Editable text outputs (.txt, .json)\n• Metrics evaluation (CER/WER)"
    p_in_bullets.font.size = Pt(11)
    p_in_bullets.font.color.rgb = c_white

    # Right Scope Box
    out_shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(4.5), Inches(5.6), Inches(2.2))
    out_shape.fill.solid()
    out_shape.fill.fore_color.rgb = c_card
    out_shape.line.color.rgb = c_red
    
    out_box = slide3.shapes.add_textbox(Inches(7.0), Inches(4.6), Inches(5.4), Inches(2.0))
    tf_out = out_box.text_frame
    tf_out.word_wrap = True
    p_out_title = tf_out.paragraphs[0]
    p_out_title.text = "OUT OF SCOPE — THIS TRIMESTER"
    p_out_title.font.size = Pt(13)
    p_out_title.font.bold = True
    p_out_title.font.color.rgb = c_red
    
    p_out_bullets = tf_out.add_paragraph()
    p_out_bullets.text = "• Full commercial web deployment\n• Perfect handwriting recognition accuracy (100%)\n• Large-scale mobile deployment\n• Complex document structure/layout understanding"
    p_out_bullets.font.size = Pt(11)
    p_out_bullets.font.color.rgb = c_white

    # =========================================================================
    # SLIDE 4: EXISTING WORK / LITERATURE
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_bg(slide4)
    add_header(slide4, "What Already Exists?", "Literature Review")
    
    # 3 horizontal nodes for timeline
    node_w = Inches(3.7)
    node_h = Inches(3.0)
    node_y = Inches(2.0)
    
    # Node 1
    n1_shape = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), node_y, node_w, node_h)
    n1_shape.fill.solid()
    n1_shape.fill.fore_color.rgb = c_card
    n1_shape.line.color.rgb = c_secondary
    
    n1_box = slide4.shapes.add_textbox(Inches(0.9), node_y + Inches(0.1), node_w - Inches(0.2), node_h - Inches(0.2))
    tf_n1 = n1_box.text_frame
    tf_n1.word_wrap = True
    p_n1 = tf_n1.paragraphs[0]
    p_n1.text = "01  Traditional HTR\n"
    p_n1.font.size = Pt(16)
    p_n1.font.bold = True
    p_n1.font.color.rgb = c_white
    
    p_n1_body = tf_n1.add_paragraph()
    p_n1_body.text = "• Base CNN + RNN + CTC loss\n• CNN for feature maps\n• LSTM/GRU for sequence modeling\n• Core pipeline: CNN → LSTM → Text"
    p_n1_body.font.size = Pt(11)
    p_n1_body.font.color.rgb = c_secondary

    # Node 2
    n2_shape = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), node_y, node_w, node_h)
    n2_shape.fill.solid()
    n2_shape.fill.fore_color.rgb = c_card
    n2_shape.line.color.rgb = c_secondary
    
    n2_box = slide4.shapes.add_textbox(Inches(4.9), node_y + Inches(0.1), node_w - Inches(0.2), node_h - Inches(0.2))
    tf_n2 = n2_box.text_frame
    tf_n2.word_wrap = True
    p_n2 = tf_n2.paragraphs[0]
    p_n2.text = "02  Transformer OCR\n"
    p_n2.font.size = Pt(16)
    p_n2.font.bold = True
    p_n2.font.color.rgb = c_white
    
    p_n2_body = tf_n2.add_paragraph()
    p_n2_body.text = "• Transformer-based sequence learning\n• Attention-based decoding\n• Removes CTC constraints\n• Core pipeline: Image → Transformer → Text"
    p_n2_body.font.size = Pt(11)
    p_n2_body.font.color.rgb = c_secondary

    # Node 3 (Highlighted)
    n3_shape = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), node_y, node_w, node_h)
    n3_shape.fill.solid()
    n3_shape.fill.fore_color.rgb = c_card
    n3_shape.line.color.rgb = c_purple
    
    n3_box = slide4.shapes.add_textbox(Inches(8.9), node_y + Inches(0.1), node_w - Inches(0.2), node_h - Inches(0.2))
    tf_n3 = n3_box.text_frame
    tf_n3.word_wrap = True
    p_n3 = tf_n3.paragraphs[0]
    p_n3.text = "03  GraDeT-HTR (2025)\n"
    p_n3.font.size = Pt(16)
    p_n3.font.bold = True
    p_n3.font.color.rgb = c_cyan
    
    p_n3_body = tf_n3.add_paragraph()
    p_n3_body.text = "• Decoder-only Transformer for Bangla HTR\n• Grapheme-based tokenization\n• Benchmark results on BN-HTRd:\n  87M Parameters\n  CER 6.19%"
    p_n3_body.font.size = Pt(11)
    p_n3_body.font.color.rgb = c_white

    # Bottom Direction
    bot_dir = slide4.shapes.add_textbox(Inches(0.8), Inches(5.3), Inches(11.7), Inches(0.8))
    tf_bd = bot_dir.text_frame
    tf_bd.word_wrap = True
    p_bd = tf_bd.paragraphs[0]
    p_bd.text = "Our Direction: Extend multilingual recognition beyond Bengali-only HTR.\nUnified HTR (Bangla + English + Digits + Symbols)"
    p_bd.font.size = Pt(13)
    p_bd.font.bold = True
    p_bd.font.color.rgb = c_cyan
    
    # Citation
    p_cit = tf_bd.add_paragraph()
    p_cit.text = "* Citation: GraDeT-HTR (2025) Bengali Handwriting Recognition. Parameters and CER cited from original publication benchmarks."
    p_cit.font.size = Pt(9)
    p_cit.font.color.rgb = c_secondary

    # =========================================================================
    # SLIDE 5: PROPOSED APPROACH
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_bg(slide5)
    add_header(slide5, "Proposed Recognition Pipeline", "Methodology")
    
    # Pipeline elements (horizontal flow)
    stages = [
        ("Input Image", "Raw source scan"),
        ("Preprocessing", "Binarize & clean"),
        ("Text Detection", "Segment lines"),
        ("Features", "CNN backbone"),
        ("Model", "BiLSTM/Transformer"),
        ("Tokenization", "Character mapping"),
        ("Digital Text", "Editable output")
    ]
    
    start_x = Inches(0.8)
    spacing = Inches(1.7)
    item_w = Inches(1.4)
    item_h = Inches(2.2)
    
    for i, (title, desc) in enumerate(stages):
        x = start_x + (i * spacing)
        
        # Draw box
        box_shape = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.3), item_w, item_h)
        box_shape.fill.solid()
        if i == 4: # Model highlight
            box_shape.fill.fore_color.rgb = c_card
            box_shape.line.color.rgb = c_purple
        else:
            box_shape.fill.fore_color.rgb = c_card
            box_shape.line.color.rgb = c_blue
            
        box_tb = slide5.shapes.add_textbox(x, Inches(2.4), item_w, item_h - Inches(0.2))
        tf_b = box_tb.text_frame
        tf_b.word_wrap = True
        
        p_t = tf_b.paragraphs[0]
        p_t.text = title
        p_t.alignment = PP_ALIGN.CENTER
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = c_cyan if i == 4 else c_white
        
        p_d = tf_b.add_paragraph()
        p_d.text = "\n" + desc
        p_d.alignment = PP_ALIGN.CENTER
        p_d.font.size = Pt(9)
        p_d.font.color.rgb = c_secondary
        
        if i < len(stages) - 1:
            # Draw an arrow
            arrow_x = x + item_w
            arrow = slide5.shapes.add_textbox(arrow_x, Inches(3.0), spacing - item_w, Inches(0.5))
            tf_a = arrow.text_frame
            p_a = tf_a.paragraphs[0]
            p_a.text = "→"
            p_a.alignment = PP_ALIGN.CENTER
            p_a.font.size = Pt(16)
            p_a.font.color.rgb = c_cyan

    # Architecture Model comparison notes at bottom
    note_shape = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.9), Inches(11.733), Inches(1.8))
    note_shape.fill.solid()
    note_shape.fill.fore_color.rgb = c_card
    note_shape.line.color.rgb = RGBColor(30, 41, 59)
    
    note_tb = slide5.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(11.3), Inches(1.6))
    tf_note = note_tb.text_frame
    tf_note.word_wrap = True
    
    p_n1 = tf_note.paragraphs[0]
    p_n1.text = "Experimental Baseline vs. Advanced Architectures:"
    p_n1.font.size = Pt(14)
    p_n1.font.bold = True
    p_n1.font.color.rgb = c_white
    
    p_n2 = tf_note.add_paragraph()
    p_n2.text = "• Baseline Model: CNN + BiLSTM + CTC (efficient, well-documented sequential approach)\n• Advanced Model: Attention-based Transformer Decoder (handling global context across languages)\n• Goal: Programmatically benchmark both architectures to determine accuracy/latency trade-offs."
    p_n2.font.size = Pt(12)
    p_n2.font.color.rgb = c_secondary

    # =========================================================================
    # SLIDE 6: DATA + TECHNOLOGY
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_bg(slide6)
    add_header(slide6, "Data & Technology Stack", "Resources & Tech")
    
    # Left Card: Datasets
    ds_shape = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5))
    ds_shape.fill.solid()
    ds_shape.fill.fore_color.rgb = c_card
    ds_shape.line.color.rgb = c_cyan
    
    ds_tb = slide6.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.1))
    tf_ds = ds_tb.text_frame
    tf_ds.word_wrap = True
    
    p_ds_t = tf_ds.paragraphs[0]
    p_ds_t.text = "📊 Training Datasets"
    p_ds_t.font.size = Pt(18)
    p_ds_t.font.bold = True
    p_ds_t.font.color.rgb = c_cyan
    
    p_ds_list = tf_ds.add_paragraph()
    p_ds_list.text = "\n• Bangla: BN-HTRd & Ishra dataset variants\n• English: IAM Handwriting Database\n• Digits: MNIST / NumtaDB digits\n• Symbols: Custom mathematical and special symbols dataset\n\n• Augmentation Strategy: Elastic distortions, shearing, rotations, blur, and contrast changes to simulate diverse handwriting profiles."
    p_ds_list.font.size = Pt(12)
    p_ds_list.font.color.rgb = c_white

    # Right Card: Technology
    tech_shape = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.5))
    tech_shape.fill.solid()
    tech_shape.fill.fore_color.rgb = c_card
    tech_shape.line.color.rgb = c_purple
    
    tech_tb = slide6.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.1))
    tf_tech = tech_tb.text_frame
    tf_tech.word_wrap = True
    
    p_tech_t = tf_tech.paragraphs[0]
    p_tech_t.text = "💻 Technical Environment"
    p_tech_t.font.size = Pt(18)
    p_tech_t.font.bold = True
    p_tech_t.font.color.rgb = c_purple
    
    p_tech_list = tf_tech.add_paragraph()
    p_tech_list.text = "\n• Programming Language: Python\n• Deep Learning: PyTorch\n• Image Processing: OpenCV, NumPy\n• API Framework: FastAPI / Flask\n• Web Interface: React JS / Vanilla HTML5\n\nNo description needed for tools - unified stack for model execution and inference web endpoint."
    p_tech_list.font.size = Pt(12)
    p_tech_list.font.color.rgb = c_white

    # =========================================================================
    # SLIDE 7: DELIVERABLES & EVALUATION
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_bg(slide7)
    add_header(slide7, "How Will We Know It Works?", "Deliverables & Metrics")
    
    # Left Card: Deliverables
    del_shape = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5))
    del_shape.fill.solid()
    del_shape.fill.fore_color.rgb = c_card
    del_shape.line.color.rgb = c_blue
    
    del_tb = slide7.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.1))
    tf_del = del_tb.text_frame
    tf_del.word_wrap = True
    
    p_del_t = tf_del.paragraphs[0]
    p_del_t.text = "📦 Project Deliverables"
    p_del_t.font.size = Pt(18)
    p_del_t.font.bold = True
    p_del_t.font.color.rgb = c_blue
    
    p_del_items = tf_del.add_paragraph()
    p_del_items.text = "\n• Dataset: Preprocessed, cleaned multilingual dataset\n• Trained Model: Weights saved (.pt/.onnx format)\n• Web Application: User interface for image uploads\n• Research Report: Project paper detailing methodologies"
    p_del_items.font.size = Pt(12)
    p_del_items.font.color.rgb = c_white

    # Right Card: Metrics Table & Cards
    met_shape = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.5))
    met_shape.fill.solid()
    met_shape.fill.fore_color.rgb = c_card
    met_shape.line.color.rgb = c_cyan
    
    met_tb = slide7.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.1))
    tf_met = met_tb.text_frame
    tf_met.word_wrap = True
    
    p_met_t = tf_met.paragraphs[0]
    p_met_t.text = "📊 Metrics & Benchmarks"
    p_met_t.font.size = Pt(18)
    p_met_t.font.bold = True
    p_met_t.font.color.rgb = c_cyan
    
    p_met_items = tf_met.add_paragraph()
    p_met_items.text = "Evaluation Criteria:\n- CER (Character Error Rate)  - WER (Word Error Rate)\n- Accuracy (Sequence-level)   - Latency (Inference speed)\n"
    p_met_items.font.size = Pt(11)
    p_met_items.font.color.rgb = c_white
    
    p_table = tf_met.add_paragraph()
    p_table.text = "Target Evaluation Schema:\n• Bangla  →  CER / Sequence Accuracy\n• English →  CER / Sequence Accuracy\n• Digits  →  Recognition Accuracy\n• Symbols →  Recognition Accuracy\n• Mixed   →  CER / WER / Edit Distance"
    p_table.font.size = Pt(11)
    p_table.font.color.rgb = c_secondary
    p_table.font.name = 'Consolas'

    # =========================================================================
    # SLIDE 8: USE OF AI ASSISTANTS
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_bg(slide8)
    add_header(slide8, "How We Use AI Assistants", "Academic Integrity")
    
    # Left Box: Workflow loop
    wf_shape = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5))
    wf_shape.fill.solid()
    wf_shape.fill.fore_color.rgb = c_card
    wf_shape.line.color.rgb = c_blue
    
    wf_tb = slide8.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.1))
    tf_wf = wf_tb.text_frame
    tf_wf.word_wrap = True
    
    p_wf_t = tf_wf.paragraphs[0]
    p_wf_t.text = "🔄 Human-AI Workflow Loop"
    p_wf_t.font.size = Pt(18)
    p_wf_t.font.bold = True
    p_wf_t.font.color.rgb = c_blue
    
    p_wf_steps = tf_wf.add_paragraph()
    p_wf_steps.text = "\nResearch\n   ↓\nAI Assistance\n   ↓\nHuman Verification & Code Audit\n   ↓\nLocal Model Implementation\n   ↓\nTesting & Experimentation"
    p_wf_steps.alignment = PP_ALIGN.CENTER
    p_wf_steps.font.size = Pt(13)
    p_wf_steps.font.bold = True
    p_wf_steps.font.color.rgb = c_white

    # Right Box: Assisting Areas & Highlight
    ast_shape = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.5))
    ast_shape.fill.solid()
    ast_shape.fill.fore_color.rgb = c_card
    ast_shape.line.color.rgb = c_purple
    
    ast_tb = slide8.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.1))
    tf_ast = ast_tb.text_frame
    tf_ast.word_wrap = True
    
    p_ast_t = tf_ast.paragraphs[0]
    p_ast_t.text = "💡 Assistive Scope"
    p_ast_t.font.size = Pt(18)
    p_ast_t.font.bold = True
    p_ast_t.font.color.rgb = c_purple
    
    p_ast_list = tf_ast.add_paragraph()
    p_ast_list.text = "• Literature: Architecture parsing & comparisons\n• Development: Debugging code, boilerplate scripts\n• Analysis: Structuring verification code\n• Writing: Presentation formatting help"
    p_ast_list.font.size = Pt(12)
    p_ast_list.font.color.rgb = c_white
    
    # Highlight Box
    p_hl = tf_ast.add_paragraph()
    p_hl.text = "\nAcademic Integrity Statement:\n\"AI assists the development process — final decisions, implementation, experiments, and validation remain our responsibility.\""
    p_hl.font.size = Pt(12)
    p_hl.font.bold = True
    p_hl.font.color.rgb = c_purple

    # =========================================================================
    # SLIDE 9: VISION (FINAL SLIDE)
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_bg(slide9)
    
    # Large center layout
    vision_tb = slide9.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(5.0))
    tf_v = vision_tb.text_frame
    tf_v.word_wrap = True
    
    p_v_tag = tf_v.paragraphs[0]
    p_v_tag.text = "LOOKING FORWARD"
    p_v_tag.alignment = PP_ALIGN.CENTER
    p_v_tag.font.size = Pt(12)
    p_v_tag.font.bold = True
    p_v_tag.font.color.rgb = c_purple
    
    p_v_title = tf_v.add_paragraph()
    p_v_title.text = "\nFrom Handwriting to Digital Intelligence"
    p_v_title.alignment = PP_ALIGN.CENTER
    p_v_title.font.size = Pt(36)
    p_v_title.font.bold = True
    p_v_title.font.color.rgb = c_white
    
    p_v_flow = tf_v.add_paragraph()
    p_v_flow.text = "\n✍️ Handwritten Page  →  👁️ AI Vision  →  🧠 HTR Model  →  💻 Editable Text"
    p_v_flow.alignment = PP_ALIGN.CENTER
    p_v_flow.font.size = Pt(16)
    p_v_flow.font.color.rgb = c_cyan
    
    p_v_goal = tf_v.add_paragraph()
    p_v_goal.text = "\n\nOne system for multilingual handwritten text.\nBangla • English • Digits • Symbols"
    p_v_goal.alignment = PP_ALIGN.CENTER
    p_v_goal.font.size = Pt(16)
    p_v_goal.font.bold = True
    p_v_goal.font.color.rgb = c_white
    
    p_ty = tf_v.add_paragraph()
    p_ty.text = "\n\nThank You!  •  Questions & Feedback"
    p_ty.alignment = PP_ALIGN.CENTER
    p_ty.font.size = Pt(16)
    p_ty.font.bold = True
    p_ty.font.color.rgb = c_green

    # Save
    prs.save("Multilingual_HTR_Presentation.pptx")
    print("Presentation saved successfully as Multilingual_HTR_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
